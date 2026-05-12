import sys
import io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

try:
    from pptx import Presentation
except ImportError:
    pass

try:
    prs = Presentation("game review.pptx")
    for i, slide in enumerate(prs.slides):
        print(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)
except Exception as e:
    print(f"Error reading pptx: {e}")
